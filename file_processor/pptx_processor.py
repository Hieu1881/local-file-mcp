from pptx import Presentation
from typing import Dict


class PPTXProcessor:
    def __init__(self, max_slides: int = 5, max_chars_per_slide: int = 2000):
        self.max_slides = max_slides
        self.max_chars_per_slide = max_chars_per_slide

    def extract_text_only(self, file_path: str) -> Dict[str, str]:
        prs = Presentation(file_path)
        text_content = {}
        slides = list(prs.slides)
        for i, slide in enumerate(slides[: self.max_slides]):
            parts = []
            for shape in slide.shapes:
                # many shapes expose .text (text frames) while others do not
                if hasattr(shape, "text"):
                    t = shape.text.strip()
                    if t:
                        parts.append(t)
            text = "\n".join(parts)
            text_content[f"slide_{i+1}"] = text[: self.max_chars_per_slide]
        return {"text_content": text_content}

    def extract_metadata_only(self, file_path: str) -> Dict[str, Dict[str, str]]:
        prs = Presentation(file_path)
        props = prs.core_properties
        metadata = {
            "title": props.title,
            "author": props.author,
            "created": str(props.created) if props.created else None,
            "last_modified_by": props.last_modified_by,
            "comments": props.comments,
            "subject": props.subject,
            "category": props.category,
        }
        return {"metadata": metadata}

    def extract_text_and_metadata(self, file_path: str) -> Dict[str, Dict[str, str]]:
        meta = self.extract_metadata_only(file_path)["metadata"]
        text = self.extract_text_only(file_path)["text_content"]
        return {"metadata": meta, "text_content": text}
